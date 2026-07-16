"""
pptx_extraction.py
-------------------
Extract structured data from .pptx files using python-pptx.

Captures per slide:
    - title / body text (all text frames, in shape order)
    - tables
    - speaker notes
    - basic key-value heuristics (e.g. "Label: Value" text boxes)

Input : ../data/pptx/*.pptx
Output:
    ../output/json/<name>_pptx.json
    ../output/dataframes/<name>_pptx_slide<i>_table<j>.csv
    ../output/key_value/<name>_pptx.txt

Install:
    pip install python-pptx pandas
"""

import json
from pathlib import Path

from pptx import Presentation
import pandas as pd

BASE_DIR = Path(__file__).resolve().parent.parent
INPUT_DIR = BASE_DIR / "data" / "pptx"
OUTPUT_JSON_DIR = BASE_DIR / "output" / "json"
OUTPUT_DF_DIR = BASE_DIR / "output" / "dataframes"
OUTPUT_KV_DIR = BASE_DIR / "output" / "key_value"

for d in (OUTPUT_JSON_DIR, OUTPUT_DF_DIR, OUTPUT_KV_DIR):
    d.mkdir(parents=True, exist_ok=True)


def extract_text_from_shape(shape):
    if not shape.has_text_frame:
        return None
    lines = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
    return lines


def extract_table_from_shape(shape):
    if not shape.has_table:
        return None
    table = shape.table
    rows = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def extract_key_value_pairs(all_text_lines):
    kv_pairs = {}
    for line in all_text_lines:
        if ":" in line:
            key, _, value = line.partition(":")
            key, value = key.strip(), value.strip()
            if key and value and len(key) < 60:
                kv_pairs[key] = value
    return kv_pairs


def process_pptx_file(filepath: Path):
    prs = Presentation(str(filepath))

    slides_data = []
    all_tables = []  # (slide_idx, table_idx, rows)
    all_text_lines = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_texts = []
        slide_tables = []

        for shape in slide.shapes:
            text_lines = extract_text_from_shape(shape)
            if text_lines:
                slide_texts.extend(text_lines)
                all_text_lines.extend(text_lines)

            table_rows = extract_table_from_shape(shape)
            if table_rows:
                slide_tables.append(table_rows)
                all_tables.append((slide_idx, len(slide_tables) - 1, table_rows))

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides_data.append({
            "slide_number": slide_idx + 1,
            "text": slide_texts,
            "tables": slide_tables,
            "speaker_notes": notes,
        })

    key_values = extract_key_value_pairs(all_text_lines)

    structured = {
        "source_file": filepath.name,
        "num_slides": len(slides_data),
        "slides": slides_data,
        "key_value_pairs": key_values,
    }

    # ---- Save JSON ----
    json_path = OUTPUT_JSON_DIR / f"{filepath.stem}_pptx.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(structured, f, indent=2, ensure_ascii=False)

    # ---- Save DataFrame(s) for each table ----
    for slide_idx, table_idx, rows in all_tables:
        if not rows:
            continue
        df = pd.DataFrame(rows[1:], columns=rows[0]) if len(rows) > 1 else pd.DataFrame(rows)
        df_path = OUTPUT_DF_DIR / f"{filepath.stem}_pptx_slide{slide_idx}_table{table_idx}.csv"
        df.to_csv(df_path, index=False)

    # ---- Save key-value text ----
    kv_path = OUTPUT_KV_DIR / f"{filepath.stem}_pptx.txt"
    with open(kv_path, "w", encoding="utf-8") as f:
        for k, v in key_values.items():
            f.write(f"{k}: {v}\n")

    print(f"[OK] Processed {filepath.name} -> {json_path.name}")
    return structured


def main():
    pptx_files = list(INPUT_DIR.glob("*.pptx"))
    if not pptx_files:
        print(f"No .pptx files found in {INPUT_DIR}")
        return

    for filepath in pptx_files:
        try:
            process_pptx_file(filepath)
        except Exception as e:
            print(f"[ERROR] Failed to process {filepath.name}: {e}")


if __name__ == "__main__":
    main()