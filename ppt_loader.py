import sys
import os
from pptx import Presentation

# Reconfigure stdout to utf-8 to handle Unicode characters/ligatures cleanly in Windows CMD
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

def find_ppt_path():
    possible_paths = ["sample.pptx", "data/sample.pptx"]
    for path in possible_paths:
        if os.path.exists(path):
            return path
    return "sample.pptx"

PPT_PATH = find_ppt_path()

def load_ppt(ppt_path=None):
    if ppt_path is None:
        ppt_path = find_ppt_path()
        
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found at '{ppt_path}'. Please place 'sample.pptx' in the project folder.")
        
    presentation = Presentation(ppt_path)
    slides_data = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_text.append(text)
        slides_data.append({
            "slide": slide_number,
            "text": "\n".join(slide_text)
        })
    return slides_data

if __name__ == "__main__":
    try:
        slides = load_ppt()
        print(f"Loaded PowerPoint file: {PPT_PATH}")
        print(f"Total Slides Parsed: {len(slides)}\n")
        for slide in slides:
            print("=" * 50)
            print(f"Slide {slide['slide']}")
            print(slide["text"].encode('ascii', errors='replace').decode('ascii'))
    except Exception as e:
        print(f"Error loading PPT: {e}")
